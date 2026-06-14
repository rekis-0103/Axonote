from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


def extract_text(file_path: Path, extension: str) -> str:
    if extension == "pdf":
        return _extract_pdf(file_path)
    if extension == "docx":
        return _extract_docx(file_path)
    if extension == "pptx":
        return _extract_pptx(file_path)
    raise ValueError(f"Unsupported extension: {extension}")


def _extract_pdf(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _extract_docx(file_path: Path) -> str:
    document = Document(str(file_path))
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())


def _extract_pptx(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)
